import os
import csv
import customtkinter as ctk
from tkinter import filedialog, messagebox
from pypdf import PdfWriter, PdfReader
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
from pdf2docx import Converter
from docx import Document
from pptx import Presentation
import docx2txt

# Set Tesseract OCR Path (Uncomment & Set If Needed)
# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# Set CustomTkinter Theme
ctk.set_appearance_mode("dark")
ctk.set_default_color_theme("blue")

# CSV File to Store User Details
USER_CSV = "users.csv"

# Create CSV File if It Doesn't Exist
if not os.path.exists(USER_CSV):
    with open(USER_CSV, mode="w", newline="") as file:
        writer = csv.writer(file)
        writer.writerow(["username", "password"])

# -------------------------------- LOGIN AND REGISTRATION FUNCTIONS --------------------------------
def register_user():
    username = register_username_entry.get()
    password = register_password_entry.get()

    if not username or not password:
        messagebox.showerror("Error", "Please enter both username and password!")
        return

    # Check if the username already exists
    with open(USER_CSV, mode="r") as file:
        reader = csv.reader(file)
        for row in reader:
            if row[0] == username:
                messagebox.showerror("Error", "Username already exists!")
                return

    # Save the new user to the CSV file
    with open(USER_CSV, mode="a", newline="") as file:
        writer = csv.writer(file)
        writer.writerow([username, password])

    messagebox.showinfo("Success", "Registration successful! You can now log in.")
    register_username_entry.delete(0, "end")
    register_password_entry.delete(0, "end")
    show_login_page()

def login_user():
    username = login_username_entry.get()
    password = login_password_entry.get()

    if not username or not password:
        messagebox.showerror("Error", "Please enter both username and password!")
        return

    # Check if the username and password match
    with open(USER_CSV, mode="r") as file:
        reader = csv.reader(file)
        for row in reader:
            if row[0] == username and row[1] == password:
                messagebox.showinfo("Success", f"Welcome, {username}!")
                login_username_entry.delete(0, "end")
                login_password_entry.delete(0, "end")
                show_main_interface()
                return

    messagebox.showerror("Error", "Invalid username or password!")

def show_register_page():
    login_frame.pack_forget()
    register_frame.pack(fill="both", expand=True)

def show_login_page():
    register_frame.pack_forget()
    login_frame.pack(fill="both", expand=True)

def show_main_interface():
    login_frame.pack_forget()
    register_frame.pack_forget()
    main_frame.pack(fill="both", expand=True)

# -------------------------------- MAIN APPLICATION FUNCTIONS --------------------------------
# Create Main Window
root = ctk.CTk()
root.title("Pdf Buddy")
root.geometry("1200x800")
root.minsize(1000, 700)

# Login Frame
login_frame = ctk.CTkFrame(root)
login_frame.pack(fill="both", expand=True)

ctk.CTkLabel(login_frame, text="🔑 Login", font=("Helvetica", 24, "bold")).pack(pady=20)
login_username_entry = ctk.CTkEntry(login_frame, placeholder_text="Username", width=300, height=40)
login_username_entry.pack(pady=10)
login_password_entry = ctk.CTkEntry(login_frame, placeholder_text="Password", show="*", width=300, height=40)
login_password_entry.pack(pady=10)
ctk.CTkButton(login_frame, text="Login", command=login_user, fg_color="#4CAF50", hover_color="#45a049", width=300, height=40).pack(pady=10)
ctk.CTkButton(login_frame, text="Register", command=show_register_page, fg_color="#008CBA", hover_color="#007B9E", width=300, height=40).pack(pady=10)

# Registration Frame
register_frame = ctk.CTkFrame(root)

ctk.CTkLabel(register_frame, text="📝 Register", font=("Helvetica", 24, "bold")).pack(pady=20)
register_username_entry = ctk.CTkEntry(register_frame, placeholder_text="Username", width=300, height=40)
register_username_entry.pack(pady=10)
register_password_entry = ctk.CTkEntry(register_frame, placeholder_text="Password", show="*", width=300, height=40)
register_password_entry.pack(pady=10)
ctk.CTkButton(register_frame, text="Register", command=register_user, fg_color="#4CAF50", hover_color="#45a049", width=300, height=40).pack(pady=10)
ctk.CTkButton(register_frame, text="Back to Login", command=show_login_page, fg_color="#008CBA", hover_color="#007B9E", width=300, height=40).pack(pady=10)

# Main Frame (Document Utility Tool)
main_frame = ctk.CTkFrame(root)

# Tab View for Navigation
tabview = ctk.CTkTabview(main_frame, width=1100, height=700, corner_radius=20)
tabview.pack(expand=True, fill="both", padx=30, pady=30)

merge_tab = tabview.add("📂 Merge PDFs")
split_tab = tabview.add("✂️ Split PDFs")
encrypt_tab = tabview.add("🔐 Encrypt PDF")
ocr_tab = tabview.add("🔍 Extract Text (OCR)")
convert_tab = tabview.add("🔄 Convert Files")

# --------------------------------- MERGE PDFs ---------------------------------
merge_file1 = ctk.StringVar()
merge_file2 = ctk.StringVar()

def select_merge_file(var):
    var.set(filedialog.askopenfilename(filetypes=[("PDF Files", "*.pdf")]))

def merge_pdfs():
    if not merge_file1.get() or not merge_file2.get():
        messagebox.showerror("Error", "Please select both PDF files!")
        return

    save_path = filedialog.asksaveasfilename(defaultextension=".pdf", title="Save Merged PDF")
    if not save_path:
        return

    writer = PdfWriter()
    writer.append(merge_file1.get())
    writer.append(merge_file2.get())

    with open(save_path, "wb") as output_pdf:
        writer.write(output_pdf)

    messagebox.showinfo("Success", "PDFs Merged Successfully! 🎉")

ctk.CTkLabel(merge_tab, text="🔗 Merge PDFs", font=("Helvetica", 24, "bold")).pack(pady=20)
ctk.CTkButton(merge_tab, text="📂 Select First PDF", command=lambda: select_merge_file(merge_file1), fg_color="#4CAF50", hover_color="#45a049").pack(pady=10)
ctk.CTkLabel(merge_tab, textvariable=merge_file1, font=("Helvetica", 12)).pack()
ctk.CTkButton(merge_tab, text="📂 Select Second PDF", command=lambda: select_merge_file(merge_file2), fg_color="#4CAF50", hover_color="#45a049").pack(pady=10)
ctk.CTkLabel(merge_tab, textvariable=merge_file2, font=("Helvetica", 12)).pack()
ctk.CTkButton(merge_tab, text="✅ Merge PDFs", fg_color="#008CBA", hover_color="#007B9E", command=merge_pdfs).pack(pady=20)

# --------------------------------- SPLIT PDFs ---------------------------------
split_file = ctk.StringVar()

def select_split_file():
    split_file.set(filedialog.askopenfilename(filetypes=[("PDF Files", "*.pdf")]))

def split_pdf():
    if not split_file.get():
        messagebox.showerror("Error", "Please select a PDF file!")
        return

    save_path1 = filedialog.asksaveasfilename(defaultextension=".pdf", title="Save First Split PDF")
    save_path2 = filedialog.asksaveasfilename(defaultextension=".pdf", title="Save Second Split PDF")

    reader = PdfReader(split_file.get())
    writer1 = PdfWriter()
    writer2 = PdfWriter()

    for i, page in enumerate(reader.pages):
        if i % 2 == 0:
            writer1.add_page(page)
        else:
            writer2.add_page(page)

    with open(save_path1, "wb") as f1, open(save_path2, "wb") as f2:
        writer1.write(f1)
        writer2.write(f2)

    messagebox.showinfo("Success", "PDF Split Successfully! 🎉")

ctk.CTkLabel(split_tab, text="✂️ Split PDFs", font=("Helvetica", 24, "bold")).pack(pady=10)
ctk.CTkButton(split_tab, text="📂 Select PDF", command=select_split_file, fg_color="#FF9800", hover_color="#FB8C00").pack(pady=5)
ctk.CTkLabel(split_tab, textvariable=split_file, font=("Helvetica", 12)).pack()
ctk.CTkButton(split_tab, text="✂️ Split PDF", fg_color="#FF5722", hover_color="#E64A19", command=split_pdf).pack(pady=10)

# --------------------------------- ENCRYPT PDF ---------------------------------
encrypt_file = ctk.StringVar()
password_entry = ctk.CTkEntry(encrypt_tab, placeholder_text="🔑 Enter Password", show="*", font=("Helvetica", 12))

def select_encrypt_file():
    encrypt_file.set(filedialog.askopenfilename(filetypes=[("PDF Files", "*.pdf")]))

def encrypt_pdf():
    if not encrypt_file.get() or not password_entry.get():
        messagebox.showerror("Error", "Please select a file and enter a password!")
        return

    save_path = filedialog.asksaveasfilename(defaultextension=".pdf", title="Save Encrypted PDF")
    if not save_path:
        return

    reader = PdfReader(encrypt_file.get())
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password_entry.get())

    with open(save_path, "wb") as f:
        writer.write(f)

    messagebox.showinfo("Success", "🔒 PDF Encrypted!")

ctk.CTkLabel(encrypt_tab, text="🔐 Encrypt PDF", font=("Helvetica", 24, "bold")).pack(pady=10)
ctk.CTkButton(encrypt_tab, text="📂 Select PDF", command=select_encrypt_file, fg_color="#9C27B0", hover_color="#8E24AA").pack(pady=5)
ctk.CTkLabel(encrypt_tab, textvariable=encrypt_file, font=("Helvetica", 12)).pack()
password_entry.pack(pady=5)
ctk.CTkButton(encrypt_tab, text="🔒 Encrypt PDF", fg_color="#E91E63", hover_color="#D81B60", command=encrypt_pdf).pack(pady=10)

# --------------------------------- EXTRACT OCR TEXT ---------------------------------
ocr_file = ctk.StringVar()

def select_ocr_file():
    ocr_file.set(filedialog.askopenfilename(filetypes=[("PDF Files", "*.pdf")]))

def extract_text_from_pdf():
    if not ocr_file.get():
        messagebox.showerror("Error", "Please select a PDF file!")
        return

    images = convert_from_path(ocr_file.get())
    extracted_text = ""

    for img in images:
        text = pytesseract.image_to_string(img)
        extracted_text += text + "\n"

    output_txt = ocr_file.get().replace(".pdf", "_text.txt")
    with open(output_txt, "w", encoding="utf-8") as f:
        f.write(extracted_text)

    messagebox.showinfo("Success", "📜 Text extracted successfully!")

ctk.CTkLabel(ocr_tab, text="📜 Extract Text (OCR)", font=("Helvetica", 24, "bold")).pack(pady=10)
ctk.CTkButton(ocr_tab, text="📂 Select PDF", command=select_ocr_file, fg_color="#673AB7", hover_color="#5E35B1").pack(pady=5)
ctk.CTkLabel(ocr_tab, textvariable=ocr_file, font=("Helvetica", 12)).pack()
ctk.CTkButton(ocr_tab, text="🔍 Extract Text", fg_color="#3F51B5", hover_color="#3949AB", command=extract_text_from_pdf).pack(pady=10)

# --------------------------------- CONVERT FILES ---------------------------------
convert_file = ctk.StringVar()
convert_format = ctk.StringVar(value="PDF to Word")

def select_convert_file():
    convert_file.set(filedialog.askopenfilename(filetypes=[("All Files", "*.*")]))

def convert_file_format():
    if not convert_file.get():
        messagebox.showerror("Error", "Please select a file!")
        return

    file_path = convert_file.get()
    format_type = convert_format.get()

    if format_type == "PDF to Word":
        save_path = filedialog.asksaveasfilename(defaultextension=".docx", title="Save As Word Document")
        if not save_path:
            return
        cv = Converter(file_path)
        cv.convert(save_path)
        cv.close()
        messagebox.showinfo("Success", "PDF converted to Word successfully! 🎉")

    elif format_type == "Word to PDF":
        save_path = filedialog.asksaveasfilename(defaultextension=".pdf", title="Save As PDF")
        if not save_path:
            return
        doc = Document(file_path)
        doc.save(save_path)
        messagebox.showinfo("Success", "Word converted to PDF successfully! 🎉")

    elif format_type == "PPT to Word":
        save_path = filedialog.asksaveasfilename(defaultextension=".docx", title="Save As Word Document")
        if not save_path:
            return
        prs = Presentation(file_path)
        doc = Document()
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    doc.add_paragraph(shape.text)
        doc.save(save_path)
        messagebox.showinfo("Success", "PPT converted to Word successfully! 🎉")

    elif format_type == "PPT to PDF":
        save_path = filedialog.asksaveasfilename(defaultextension=".pdf", title="Save As PDF")
        if not save_path:
            return
        prs = Presentation(file_path)
        prs.save(save_path)
        messagebox.showinfo("Success", "PPT converted to PDF successfully! 🎉")

    elif format_type == "Word to TXT":
        save_path = filedialog.asksaveasfilename(defaultextension=".txt", title="Save As Text File")
        if not save_path:
            return
        text = docx2txt.process(file_path)
        with open(save_path, "w", encoding="utf-8") as f:
            f.write(text)
        messagebox.showinfo("Success", "Word converted to TXT successfully! 🎉")

    elif format_type == "PDF to TXT":
        save_path = filedialog.asksaveasfilename(defaultextension=".txt", title="Save As Text File")
        if not save_path:
            return
        images = convert_from_path(file_path)
        extracted_text = ""
        for img in images:
            text = pytesseract.image_to_string(img)
            extracted_text += text + "\n"
        with open(save_path, "w", encoding="utf-8") as f:
            f.write(extracted_text)
        messagebox.showinfo("Success", "PDF converted to TXT successfully! 🎉")

ctk.CTkLabel(convert_tab, text="🔄 Convert Files", font=("Helvetica", 24, "bold")).pack(pady=10)
ctk.CTkButton(convert_tab, text="📂 Select File", command=select_convert_file, fg_color="#009688", hover_color="#00897B").pack(pady=5)
ctk.CTkLabel(convert_tab, textvariable=convert_file, font=("Helvetica", 12)).pack()
ctk.CTkOptionMenu(convert_tab, variable=convert_format, values=["PDF to Word", "Word to PDF", "PPT to Word", "PPT to PDF", "Word to TXT", "PDF to TXT"], fg_color="#00BCD4").pack(pady=5)
ctk.CTkButton(convert_tab, text="🔄 Convert", fg_color="#FFC107", hover_color="#FFB300", command=convert_file_format).pack(pady=10)

# Show Login Page by Default
show_login_page()

# Run the App
root.mainloop()